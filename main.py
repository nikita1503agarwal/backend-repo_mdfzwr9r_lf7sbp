import os
import io
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def read_root():
    return {"message": "Hello from FastAPI Backend!"}

@app.get("/api/hello")
def hello():
    return {"message": "Hello from the backend API!"}

@app.get("/test")
def test_database():
    """Test endpoint to check if database is available and accessible"""
    response = {
        "backend": "✅ Running",
        "database": "❌ Not Available",
        "database_url": None,
        "database_name": None,
        "connection_status": "Not Connected",
        "collections": []
    }
    
    try:
        # Try to import database module
        from database import db
        
        if db is not None:
            response["database"] = "✅ Available"
            response["database_url"] = "✅ Configured"
            response["database_name"] = db.name if hasattr(db, 'name') else "✅ Connected"
            response["connection_status"] = "Connected"
            
            # Try to list collections to verify connectivity
            try:
                collections = db.list_collection_names()
                response["collections"] = collections[:10]  # Show first 10 collections
                response["database"] = "✅ Connected & Working"
            except Exception as e:
                response["database"] = f"⚠️  Connected but Error: {str(e)[:50]}"
        else:
            response["database"] = "⚠️  Available but not initialized"
            
    except ImportError:
        response["database"] = "❌ Database module not found (run enable-database first)"
    except Exception as e:
        response["database"] = f"❌ Error: {str(e)[:50]}"
    
    # Check environment variables
    import os
    response["database_url"] = "✅ Set" if os.getenv("DATABASE_URL") else "❌ Not Set"
    response["database_name"] = "✅ Set" if os.getenv("DATABASE_NAME") else "❌ Not Set"
    
    return response


# --- PPTX Export Endpoint ---
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def build_deck_content():
    """Returns the structured content used by the frontend deck."""
    return [
        {"type": "title", "title": "AI Study Bot", "subtitle": "Intelligent Assistant for Learning"},
        {"title": "Introduction", "points": [
            'An AI-powered assistant that helps students learn faster and smarter',
            'Provides instant explanations, study aids, and guidance across subjects',
            'Combines NLP, retrieval, and personalization to support learning goals',
        ]},
        {"title": "Problem Statement", "points": [
            'Information overload across textbooks, notes, and the web',
            'Lack of tailored guidance for different learners',
            'Manual notes are time-consuming and hard to organize',
            'Revising efficiently is challenging without structure',
        ]},
        {"title": "Objectives", "points": [
            'Instant explanations on-demand',
            'Automatic summaries of long content',
            'Personalized learning paths',
            'Quizzes to test understanding',
            'Progress tracking and analytics',
        ]},
        {"title": "System Architecture", "points": [
            'User Interface',
            'NLP Engine',
            'Knowledge Base',
            'Recommendation Module',
            'Feedback Loop',
        ]},
        {"title": "Workflow", "flow": [
            'User query', 'NLP', 'Retrieval', 'Response', 'Personalized suggestion'
        ]},
        {"title": "Key Features", "points": [
            'Doubt solving and instant Q&A',
            'Summarization of lectures and notes',
            'Quiz generation with adaptive difficulty',
            'Progress tracking dashboards',
            'Adaptive learning recommendations',
        ]},
        {"title": "Use Cases", "points": [
            'Exam preparation',
            'Quick revision sessions',
            'Assignments and research support',
            'Self-study and continuous learning',
        ]},
        {"title": "Technology Stack", "points": [
            'Python for backend logic',
            'Modern NLP models (LLMs, embeddings)',
            'Vector databases for retrieval',
            'Recommendation algorithms',
            'Web frameworks for UI + APIs',
        ]},
        {"title": "Benefits", "points": [
            'Deeper understanding with timely help',
            'Reduced study time through automation',
            '24/7 assistance and motivation',
        ]},
        {"title": "Limitations", "points": [
            'May produce inaccuracies; verification needed',
            'Requires regular data updates',
            'Possible misunderstandings of context',
            'Internet connectivity required',
        ]},
        {"title": "Future Enhancements", "points": [
            'Voice interaction and speech synthesis',
            'Multimodal input (images, PDFs, whiteboards)',
            'Offline mode for core features',
            'Integration with LMS platforms',
        ]},
        {"title": "Conclusion", "points": [
            'AI Study Bot streamlines learning, clarifies doubts, and adapts to each student',
            'Boosts outcomes through explanations, quizzes, and progress insight',
            'A reliable companion for effective, continuous learning',
        ]},
    ]


def create_pptx_bytes():
    deck = build_deck_content()
    prs = PptxPresentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    first = deck[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = first.get("title", "AI Study Bot")
    subtitle_placeholder = slide.placeholders[1]
    subtitle_placeholder.text = first.get("subtitle", "Intelligent Assistant for Learning")

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title + Content

    for section in deck[1:]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section.get("title", "")
        tf = slide.placeholders[1].text_frame
        tf.clear()

        # Choose between points and flow
        points = section.get("points")
        flow = section.get("flow")

        if points:
            for i, p in enumerate(points):
                p_run = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p_run.text = p
                p_run.level = 0
        elif flow:
            for i, step in enumerate(flow):
                p_run = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                arrow = " → " if i < len(flow) - 1 else ""
                p_run.text = f"{step}{arrow}"
                p_run.level = 0
        else:
            tf.paragraphs[0].text = ""

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


@app.get("/api/export/pptx")
def export_pptx():
    try:
        bio = create_pptx_bytes()
        filename = "AI_Study_Bot_Presentation.pptx"
        headers = {
            "Content-Disposition": f"attachment; filename={filename}"
        }
        return StreamingResponse(bio, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
